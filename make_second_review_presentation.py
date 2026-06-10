"""Generate second_review_presentation.pptx for the 6 June 2026 review."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x2F, 0x5E)
BLUE      = RGBColor(0x1A, 0x6F, 0xC4)
SKY       = RGBColor(0xD0, 0xE8, 0xF8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL  = RGBColor(0x22, 0x22, 0x22)
MID_GRAY  = RGBColor(0x55, 0x55, 0x55)
GREEN     = RGBColor(0x0D, 0x7A, 0x3E)
ORANGE    = RGBColor(0xD6, 0x63, 0x10)


def set_cell_bg(cell, rgb: RGBColor):
    """Fill a table cell with a solid background colour."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgb.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


def add_rect(slide, l, t, w, h, fill: RGBColor, alpha=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h),
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def title_text(txBox, text, size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT):
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text_box(slide, l, t, w, h, text, size=22, bold=False,
                 color=CHARCOAL, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, l, t, w, h, items, size=22, color=CHARCOAL,
                    bullet_color=BLUE, indent=0.25):
    """Add a bulleted list; items can be str or (str, indent_level)."""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        indent_inches = indent * (level + 1)
        p.space_before = Pt(4)

        # bullet character
        bullet_run = p.add_run()
        bullet_char = "▸" if level == 0 else "–"
        bullet_run.text = bullet_char + "  "
        bullet_run.font.size = Pt(size)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = bullet_color

        content_run = p.add_run()
        content_run.text = text
        content_run.font.size = Pt(size)
        content_run.font.color.rgb = color
    return txBox


def add_header_bar(slide, title, subtitle=None):
    """Add the top navy header strip with white title text."""
    add_rect(slide, 0, 0, 13.33, 1.35, NAVY)
    add_rect(slide, 0, 1.35, 13.33, 0.07, BLUE)

    txBox = slide.shapes.add_textbox(Inches(0.45), Inches(0.1),
                                      Inches(12.5), Inches(0.85))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(38)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.45), Inches(0.88),
                                           Inches(12.5), Inches(0.45))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(20)
        run2.font.color.rgb = SKY


def add_footer(slide, slide_num, total=5):
    add_rect(slide, 0, 7.1, 13.33, 0.4, NAVY)
    add_text_box(slide, 0.3, 7.12, 8, 0.35,
                 "AI-Assisted 3D Assembly Design  |  PES University Bengaluru",
                 size=14, color=SKY)
    add_text_box(slide, 11.8, 7.12, 1.3, 0.35,
                 f"{slide_num} / {total}", size=14, color=SKY,
                 align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title and Abstract
# ═══════════════════════════════════════════════════════════════════════════════
def slide1_title_abstract(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # full-slide gradient-like background
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5, 0xF8, 0xFF))

    # left accent bar
    add_rect(slide, 0, 0, 0.18, 7.5, NAVY)

    # top decorative strip
    add_rect(slide, 0.18, 0, 13.15, 0.06, BLUE)

    # Institution tag
    add_text_box(slide, 1.0, 0.35, 12, 0.5,
                 "PES University Bengaluru  |  M.Tech Individual Project — Semester 4",
                 size=19, color=MID_GRAY, italic=True)

    # Main title
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.0),
                                      Inches(11.5), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "AI-Assisted 3D Assembly Design"
    run.font.size = Pt(52)
    run.font.bold = True
    run.font.color.rgb = NAVY

    # Subtitle
    add_text_box(slide, 1.0, 2.7, 11.5, 0.6,
                 "Second Review Presentation",
                 size=32, bold=True, color=BLUE)

    # Date line
    add_text_box(slide, 1.0, 3.25, 11.5, 0.45,
                 "Date: 6th June 2026",
                 size=22, color=MID_GRAY)

    # Divider
    add_rect(slide, 1.0, 3.85, 10.5, 0.04, BLUE)

    # Abstract heading
    add_text_box(slide, 1.0, 4.0, 11.5, 0.45,
                 "Abstract", size=26, bold=True, color=NAVY)

    abstract = (
        "This project develops an end-to-end AI pipeline that analyses partial 3D mechanical "
        "assemblies supplied as STEP files, constructs graph representations of parts and their "
        "mating relationships, and employs a Graph Attention Network (GAT) to detect missing "
        "components.  A natural-language assistant (AIDA), powered by Gemini 2.5 Flash, explains "
        "assembly structure and suggests corrective actions.  The system is delivered as an "
        "interactive web application combining a FastAPI backend with a Streamlit frontend."
    )
    add_text_box(slide, 1.0, 4.5, 11.5, 2.2,
                 abstract, size=21, color=CHARCOAL)

    # Student card at bottom-right
    add_rect(slide, 8.8, 6.4, 4.35, 0.95, SKY)
    add_text_box(slide, 8.95, 6.43, 4.1, 0.45,
                 "Presented by", size=16, color=MID_GRAY, italic=True)
    add_text_box(slide, 8.95, 6.72, 4.1, 0.55,
                 "Parthasarathy Perumal", size=21, bold=True, color=NAVY)

    # Bottom accent
    add_rect(slide, 0.18, 7.44, 13.15, 0.06, BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — Detailed Algorithm Design
# ═══════════════════════════════════════════════════════════════════════════════
def slide2_algo_design(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5, 0xF8, 0xFF))
    add_header_bar(slide, "Detailed Algorithm Design",
                   "Graph Neural Network Pipeline for Missing Component Detection")
    add_footer(slide, 2)

    # ── Stage boxes ─────────────────────────────────────────────────────────
    stages = [
        ("1  Input\nProcessing",    NAVY,  0.25),
        ("2  Graph\nConstruction",  BLUE,  2.92),
        ("3  GNN\nModel",           GREEN, 5.59),
        ("4  Prediction\n& Output", ORANGE,8.26),
    ]
    for label, color, x in stages:
        add_rect(slide, x, 1.55, 2.6, 1.15, color)
        add_text_box(slide, x + 0.07, 1.62, 2.46, 1.0,
                     label, size=21, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Arrow (except after last)
        if x < 8.26:
            add_text_box(slide, x + 2.6, 1.92, 0.34, 0.45,
                         "→", size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ── Detail bullets per stage ─────────────────────────────────────────────
    stage_details = [
        [
            "STEP file as input",
            "gmsh / OpenCASCADE parser",
            "Extracts 3D solid geometry",
            "Generates mesh & B-rep data",
        ],
        [
            "Nodes — 13-dim features:",
            ("Part-type one-hot (8 dims)", 1),
            ("Volume, Surface area", 1),
            ("Bounding-box Δx, Δy, Δz", 1),
            "Edges — 2-dim features:",
            ("Mate-type encoding", 1),
            ("Contact weight", 1),
        ],
        [
            "3-layer Graph Attention Net",
            "Multi-head attention (8 heads)",
            "64-dim node embeddings",
            "Dropout = 0.1 for regularisation",
            "Adam + ReduceLROnPlateau",
        ],
        [
            "MLP Link Predictor",
            "Binary edge score per pair",
            "Threshold → missing parts",
            "NodeRanker for ranking",
            "AIDA (Gemini) explanation",
        ],
    ]
    xs = [0.25, 2.92, 5.59, 8.26]
    for i, (details, x) in enumerate(zip(stage_details, xs)):
        y = 2.82
        for item in details:
            if isinstance(item, tuple):
                label, level = item
                add_text_box(slide, x + 0.35, y, 2.25, 0.35,
                             f"  – {label}", size=17, color=MID_GRAY)
            else:
                add_text_box(slide, x + 0.05, y, 2.55, 0.38,
                             f"▸ {item}", size=18, bold=False, color=CHARCOAL)
            y += 0.36

    # ── AIDA / Skills block ──────────────────────────────────────────────────
    add_rect(slide, 0.25, 6.1, 12.83, 0.82, SKY)
    add_text_box(slide, 0.4, 6.14, 2.0, 0.38,
                 "AIDA Assistant:", size=20, bold=True, color=NAVY)
    add_text_box(slide, 2.35, 6.14, 10.6, 0.38,
                 "Gemini 2.5 Flash LLM  +  engineering_3d_assembly persona  →  "
                 "Natural-language explanations, assembly-category analysis, "
                 "joint-type identification, and corrective suggestions via 6 domain skills.",
                 size=19, color=CHARCOAL)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — Student Contribution
# ═══════════════════════════════════════════════════════════════════════════════
def slide3_contribution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5, 0xF8, 0xFF))
    add_header_bar(slide, "Contribution of the Student",
                   "Parthasarathy Perumal — Individual Project, M.Tech Sem 4")
    add_footer(slide, 3)

    col_w = 5.9
    # ── Left column ──────────────────────────────────────────────────────────
    add_rect(slide, 0.3, 1.55, col_w, 0.48, NAVY)
    add_text_box(slide, 0.38, 1.57, col_w - 0.1, 0.44,
                 "Core AI / ML Pipeline", size=22, bold=True, color=WHITE)

    left_items = [
        "Designed & implemented the full STEP-to-GNN pipeline from scratch",
        "Built STEP parser using gmsh + OpenCASCADE — handles real CAD files",
        "Engineered 13-dimensional node feature vectors from 3D geometry",
        "Designed AssemblyGNN: 3-layer GAT with multi-head attention",
        "Implemented MLP Link Predictor for binary missing-part detection",
        "Wrote training loop with Adam, early stopping, ReduceLROnPlateau",
        "Defined evaluation metrics: AUC-ROC and Average Precision (AP)",
    ]
    y = 2.12
    for item in left_items:
        add_text_box(slide, 0.55, y, col_w - 0.3, 0.44,
                     f"▸  {item}", size=19, color=CHARCOAL)
        y += 0.48

    # ── Right column ─────────────────────────────────────────────────────────
    rx = 7.05
    add_rect(slide, rx, 1.55, col_w, 0.48, BLUE)
    add_text_box(slide, rx + 0.08, 1.57, col_w - 0.1, 0.44,
                 "System, UI & AI Skills Agent", size=22, bold=True, color=WHITE)

    right_items = [
        "Built FastAPI backend exposing /predict/missing, /explain, /identify, /health",
        "Created Streamlit frontend with Plotly 3D mesh viewer of assemblies",
        "Integrated real-time training monitor within the web dashboard",
        "Developed AIDA agent (Gemini 2.5 Flash) with 6 domain skills YAML",
        "Engineered system prompt / persona (engineering_3d_assembly.yaml)",
        "Wrote shell scripts for one-command service start/stop/log access",
        "Maintained a version-controlled project across macOS (MPS) environment",
    ]
    y = 2.12
    for item in right_items:
        add_text_box(slide, rx + 0.2, y, col_w - 0.3, 0.44,
                     f"▸  {item}", size=19, color=CHARCOAL)
        y += 0.48

    # Divider between columns
    add_rect(slide, 6.67, 1.55, 0.04, 5.35, SKY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — Results (Intermediate)
# ═══════════════════════════════════════════════════════════════════════════════
def slide4_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5, 0xF8, 0xFF))
    add_header_bar(slide, "Results Obtained  —  Intermediate",
                   "GNN Training Metrics as of May 2026  |  Phase 1 Complete")
    add_footer(slide, 4)

    # ── Metric highlight cards ────────────────────────────────────────────────
    cards = [
        ("Best AUC-ROC",   "0.9272", "assembly_gnn_20260523\n_162431", NAVY),
        ("Avg Precision",  "0.6463", "Current test set\nAP score",     BLUE),
        ("Models Trained", "9",      "Checkpoints saved\nin trained_models/", GREEN),
        ("Training Runs",  "6+",     "Multiple configs\nexplored",      ORANGE),
    ]
    for i, (label, value, note, color) in enumerate(cards):
        x = 0.3 + i * 3.18
        add_rect(slide, x, 1.58, 2.95, 1.55, color)
        add_text_box(slide, x + 0.12, 1.65, 2.7, 0.42,
                     label, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.12, 2.05, 2.7, 0.72,
                     value, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.12, 2.82, 2.7, 0.25,
                     note, size=14, color=RGBColor(0xDD, 0xEE, 0xFF),
                     align=PP_ALIGN.CENTER)

    # ── Training configuration ────────────────────────────────────────────────
    add_rect(slide, 0.3, 3.35, 5.9, 0.42, NAVY)
    add_text_box(slide, 0.4, 3.37, 5.7, 0.38,
                 "Training Configuration", size=21, bold=True, color=WHITE)

    config = [
        ("Optimiser",        "Adam  (lr = 1×10⁻³)"),
        ("LR Scheduler",     "ReduceLROnPlateau (factor 0.5, patience 10)"),
        ("Early Stopping",   "Patience = 20 epochs"),
        ("Batch strategy",   "Full-graph (single batch per epoch)"),
        ("Input features",   "Node: 13-dim  |  Edge: 2-dim"),
        ("Embedding dim",    "64  (hidden = 128)"),
        ("Dropout",          "0.1"),
        ("Device",           "MPS (Apple Silicon M-series)"),
    ]
    y = 3.83
    for key, val in config:
        add_text_box(slide, 0.45, y, 2.2, 0.38,
                     key, size=18, bold=True, color=NAVY)
        add_text_box(slide, 2.65, y, 3.35, 0.38,
                     val, size=18, color=CHARCOAL)
        y += 0.4

    # ── AUC variation note ────────────────────────────────────────────────────
    add_rect(slide, 6.55, 3.35, 6.5, 0.42, BLUE)
    add_text_box(slide, 6.65, 3.37, 6.3, 0.38,
                 "AUC Variation Across Runs", size=21, bold=True, color=WHITE)

    runs = [
        ("Run 1 (best)",  "AUC = 0.9272",  "Optimal LR + data split"),
        ("Run 2",         "AUC = 0.8641",  "Slightly larger dropout"),
        ("Run 3",         "AUC = 0.7418",  "Reduced hidden dim"),
        ("Run 4",         "AUC = 0.6367",  "Latest checkpoint"),
        ("Run 5",         "AUC = 0.4424",  "Current test_metrics.json"),
        ("Observation",   "High variance", "Sensitivity to data split & LR"),
    ]
    y = 3.83
    for run, metric, note in runs:
        col = GREEN if "best" in run else (ORANGE if float(metric.split("=")[-1].strip()) < 0.5 else CHARCOAL) if "=" in metric else BLUE
        add_text_box(slide, 6.65, y, 2.0, 0.38,
                     run, size=18, bold=True, color=NAVY)
        add_text_box(slide, 8.65, y, 1.6, 0.38,
                     metric, size=18, bold=True, color=col)
        add_text_box(slide, 10.25, y, 2.65, 0.38,
                     note, size=16, italic=True, color=MID_GRAY)
        y += 0.4

    # ── Next steps ────────────────────────────────────────────────────────────
    add_rect(slide, 0.3, 6.78, 12.73, 0.06, BLUE)
    add_text_box(slide, 0.3, 6.5, 12.5, 0.42,
                 "Next (Phase 2):  Next-component ranking using NodeRanker  |  "
                 "Metrics: Hit@K, MRR, NDCG@K  |  Target AUC ≥ 0.92 consistently",
                 size=18, color=NAVY, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — References
# ═══════════════════════════════════════════════════════════════════════════════
def slide5_references(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0xF5, 0xF8, 0xFF))
    add_header_bar(slide, "References",
                   "Key Literature, Frameworks, and Tools")
    add_footer(slide, 5)

    refs = [
        ("[1]",
         "P. Veličković et al.",
         "Graph Attention Networks.",
         "ICLR 2018.  arXiv:1710.10903"),

        ("[2]",
         "M. Fey & J. E. Lenssen.",
         "Fast Graph Representation Learning with PyTorch Geometric.",
         "ICLR Workshop on Representation Learning on Graphs and Manifolds, 2019."),

        ("[3]",
         "Geuzaine, C. & Remacle, J.-F.",
         "Gmsh: A 3-D finite element mesh generator with built-in pre- and post-processing facilities.",
         "International Journal for Numerical Methods in Engineering, 79(11), 1309–1331, 2009."),

        ("[4]",
         "Open CASCADE Technology (OCCT).",
         "Open CASCADE 3D Modeling Kernel.",
         "https://www.opencascade.com  (accessed June 2026)"),

        ("[5]",
         "Google DeepMind.",
         "Gemini 2.5 Flash — Multimodal Large Language Model.",
         "Technical Report, Google LLC, 2025."),

        ("[6]",
         "FastAPI Framework.",
         "Sebastián Ramírez.  FastAPI — High Performance, Easy to Learn, Fast to Code.",
         "https://fastapi.tiangolo.com  (accessed June 2026)"),

        ("[7]",
         "Streamlit Inc.",
         "Streamlit — A Faster Way to Build and Share Data Apps.",
         "https://streamlit.io  (accessed June 2026)"),

        ("[8]",
         "A. Paszke et al.",
         "PyTorch: An Imperative Style, High-Performance Deep Learning Library.",
         "NeurIPS 2019.  arXiv:1912.01703"),
    ]

    y = 1.6
    for num, authors, title, venue in refs:
        add_text_box(slide, 0.35, y, 0.55, 0.5,
                     num, size=20, bold=True, color=BLUE)
        add_text_box(slide, 0.9, y, 11.9, 0.22,
                     authors + "  " + title,
                     size=20, bold=False, color=CHARCOAL)
        add_text_box(slide, 0.9, y + 0.24, 11.9, 0.22,
                     venue, size=17, italic=True, color=MID_GRAY)
        y += 0.59


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide1_title_abstract(prs)
    slide2_algo_design(prs)
    slide3_contribution(prs)
    slide4_results(prs)
    slide5_references(prs)

    out = "second_review_presentation.pptx"
    prs.save(out)
    print(f"Saved → {out}")


if __name__ == "__main__":
    main()
